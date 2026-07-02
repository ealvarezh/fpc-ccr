import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# =====================================================================
# CONFIGURACION
# =====================================================================
SILVER = Path(r"C:\Users\eah\apoyoconsultoria.com\File Server - Analytics\7 Datos\Datos abiertos\fissal\01_silver")
OUTPUT_DIR = SILVER.parent / "03_output"
PPTX_OUT = OUTPUT_DIR / "FISSAL_CCR_Hitos_Presentacion.pptx"

PRIMARY   = RGBColor(0x11, 0x8D, 0xFF)
FOREGROUND= RGBColor(0x25, 0x24, 0x23)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF3, 0xF2, 0xF1)
GOOD      = RGBColor(0x1A, 0xAB, 0x40)
NEUTRAL   = RGBColor(0xD9, 0xB3, 0x00)
BAD       = RGBColor(0xD6, 0x45, 0x54)
DARK_BLUE = RGBColor(0x0D, 0x5E, 0xAA)
MID_BLUE  = RGBColor(0x43, 0x93, 0xC3)
GRAY      = RGBColor(0x60, 0x5E, 0x5C)
LIGHTER   = RGBColor(0xC0, 0xD0, 0xE8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# =====================================================================
# CARGA DE DATOS
# =====================================================================
print("Cargando datos de hitos...")
hitos = pd.read_parquet(SILVER / "FISSAL_CCR_HITOS_PACIENTE.parquet")
N_PAC = len(hitos)

# =====================================================================
# FUNCIONES
# =====================================================================
def add_slide(layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def add_title(slide, text, left=0.5, top=0.3, width=12.3, height=0.7, font_size=28):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = True; p.font.color.rgb = FOREGROUND; p.font.name = "Segoe UI"

def add_subtitle(slide, text, left=0.5, top=1.0, width=12.3, height=0.4, font_size=15):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.color.rgb = GRAY; p.font.name = "Segoe UI"

def add_text(slide, text, left, top, width, height, font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
    if color is None: color = FOREGROUND
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = "Segoe UI"; p.alignment = align

def add_bullet(slide, text, left, top, width, height, font_size=11, color=None, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color or FOREGROUND; p.font.name = "Segoe UI"

def add_kpi(slide, label, value, left, top, width=2.8, height=1.4, color=None):
    if color is None: color = PRIMARY
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE1, 0xDF, 0xDD); shape.line.width = Pt(1)
    add_text(slide, str(value), left + 0.2, top + 0.15, width - 0.4, 0.65, font_size=28, bold=True, color=color)
    add_text(slide, label, left + 0.2, top + 0.8, width - 0.4, 0.4, font_size=11, color=GRAY)

def add_table(slide, headers, rows, left, top, col_widths, font_size=10):
    n_rows = len(rows) + 1; n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(total_w), Inches(0.35 * n_rows))
    table = table_shape.table
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
        for p in table.cell(0, j).text_frame.paragraphs:
            p.font.size = Pt(font_size); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
        table.cell(0, j).fill.solid(); table.cell(0, j).fill.fore_color.rgb = PRIMARY
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            table.cell(i + 1, j).text = str(val)
            for p in table.cell(i + 1, j).text_frame.paragraphs:
                p.font.size = Pt(font_size); p.font.color.rgb = FOREGROUND; p.font.name = "Segoe UI"
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            table.cell(i + 1, j).fill.solid()
            table.cell(i + 1, j).fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_BG
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)

def add_box(slide, text, left, top, width, height=0.6):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BG; shape.line.fill.background()
    add_text(slide, text, left + 0.15, top + 0.08, width - 0.3, height - 0.15, font_size=9, color=GRAY)

def add_footer(slide):
    add_text(slide, "FISSAL CCR | Hitos del Paciente | Atenciones 2012-2022", 0.5, 7.0, 12, 0.3, font_size=8, color=RGBColor(0xA1, 0x9F, 0x9D))

def add_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = PRIMARY; shape.line.fill.background()

# =====================================================================
# SLIDE 1: PORTADA
# =====================================================================
slide = add_slide()
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE
add_text(slide, "FISSAL  |  CÁNCER COLORRECTAL", 1, 1.8, 11, 1.0, font_size=36, bold=True, color=WHITE)
add_text(slide, "Hitos del Paciente: Despistaje → Diagnóstico → Tratamiento → Desenlace", 1, 2.9, 11, 0.7, font_size=18, color=LIGHTER)
add_line(slide, 1, 3.7, 4)
add_text(slide, "Reconstrucción de la trayectoria clínica completa  |  15,698 pacientes  |  C18/C19/C20", 1, 4.1, 11, 0.5, font_size=13, color=LIGHTER)
add_text(slide, "Fuente: FISSAL — Prestaciones de salud 2012-2022 (procesadas 2016-2022)", 1, 5.0, 11, 0.5, font_size=11, color=LIGHTER)

# =====================================================================
# SLIDE 2: LOS 4 HITOS — VISIÓN GENERAL
# =====================================================================
slide = add_slide()
add_title(slide, "Modelo de 4 Hitos del Paciente")
add_subtitle(slide, "Estructura conceptual que organiza toda la trayectoria clínica del paciente con CCR")
add_footer(slide)

hitos_info = [
    ("1. DESPISTAJE", "Procedimientos de sospecha ANTES del diagnóstico formal.\nSe busca en data COMPLETA (todos los CIE-10) de cada paciente.", PRIMARY),
    ("2. DIAGNÓSTICO", "Primera atención codificada como C18/C19/C20.\nTipo de ingreso (ambulatorio/emergencia/hosp) como proxy de gravedad.", MID_BLUE),
    ("3. TRATAMIENTO", "Episodios de Cirugía, Quimioterapia y Radioterapia.\nSecuencia, sesiones, costos, hospitalizaciones y adherencia.", GOOD),
    ("4. DESENLACE", "6 categorías: Fallecido, Paliativo, Recaída probable,\nControl post-pausa, Posible remisión/alta, Seguimiento activo.", BAD),
]

for i, (tit, desc, color) in enumerate(hitos_info):
    x = 0.5 + i * 3.15
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(2.9), Inches(2.0))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(slide, tit, x + 0.2, 1.9, 2.5, 0.5, font_size=14, bold=True, color=WHITE)
    add_text(slide, desc, x + 0.2, 2.5, 2.5, 1.2, font_size=9, color=WHITE)

n_desp = hitos["TUVO_DESPISTAJE_PREVIO"].sum()
n_dx = N_PAC
n_trat = (hitos["SECUENCIA_TRATAMIENTO"] != "SIN_TRATAMIENTO_DETECTADO").sum()
n_fall = (hitos["DESENLACE"] == "FALLECIDO").sum()

add_text(slide, "Embudo de pacientes", 0.5, 4.2, 5, 0.4, font_size=14, bold=True)
embudo_rows = [
    ["Despistaje detectado en FISSAL", f"{n_desp:,}", f"{n_desp/N_PAC*100:.1f}%"],
    ["Diagnóstico CCR (base)", f"{n_dx:,}", "100.0%"],
    ["Tratamiento detectado", f"{n_trat:,}", f"{n_trat/N_PAC*100:.1f}%"],
    ["Fallecidos", f"{n_fall:,}", f"{n_fall/N_PAC*100:.1f}%"],
]
add_table(slide, ["Hito", "Pacientes", "% del total"], embudo_rows, 0.5, 4.7, [3.5, 1.5, 1.2], font_size=10)

add_box(slide, "Solo 6.6% tiene despistaje detectable en FISSAL. El 93.4% aparece por primera vez YA con el diagnóstico CCR. Esto NO significa que no hubo despistaje: probablemente ocurrió en otro subsistema (EsSalud, consulta privada, MINSA sin FISSAL).", 6.8, 4.2, 6.0, 1.0)

# =====================================================================
# SLIDE 3: HITO 1 — DESPISTAJE
# =====================================================================
slide = add_slide()
add_title(slide, "Hito 1: Despistaje / Detección")
add_subtitle(slide, "Eventos de sospecha ANTES del diagnóstico C18/C19/C20, buscados en la data completa de FISSAL (todos los diagnósticos)")
add_footer(slide)

con_desp = hitos[hitos["TUVO_DESPISTAJE_PREVIO"]]
sin_desp = hitos[~hitos["TUVO_DESPISTAJE_PREVIO"]]

add_kpi(slide, "Con despistaje detectado", f"{len(con_desp):,} ({len(con_desp)/N_PAC*100:.1f}%)", 0.5, 1.8, color=GOOD)
add_kpi(slide, "Sin despistaje (dx es 1er contacto)", f"{len(sin_desp):,} ({len(sin_desp)/N_PAC*100:.1f}%)", 3.6, 1.8, color=NEUTRAL)

# Tipo de despistaje
tipo_desp = con_desp["TIPO_DESPISTAJE"].value_counts()
rows_tipo = [[t, n, f"{n/len(con_desp)*100:.1f}%"] for t, n in tipo_desp.items()]
add_text(slide, "Tipo de primer evento de despistaje", 0.5, 3.5, 5, 0.4, font_size=13, bold=True)
add_table(slide, ["Tipo", "Pacientes", "%"], rows_tipo, 0.5, 4.0, [2.5, 1.5, 1.0])

# Categorias de búsqueda
add_text(slide, "Categorías buscadas en datos pre-diagnóstico", 6.8, 3.5, 5, 0.4, font_size=13, bold=True)
cats_desp = [
    ["BIOPSIA", "BIOPSI | ANATOMOPATOLOG"],
    ["ENDOSCOPIA", "COLONOSCOP | ENDOSCOP | RECTOSCOP | SIGMOIDOSCOP"],
    ["IMAGEN", "TOMOGRAF | RESONAN | ECOGRAF | GAMMAGRAF | PET-SCAN | RADIOGRAF"],
    ["LABORATORIO", "CEA | ANTIGEN CARCINO | SANGRE OCULTA"],
]
add_table(slide, ["Categoría", "Patrón regex"], cats_desp, 6.8, 4.0, [2.0, 3.5])

add_box(slide, "Los despistajes se buscan en los archivos anuales COMPLETOS (todos los CIE-10, no solo C18/19/20), en fechas ANTERIORES a la primera atención CCR de cada paciente. Un despistaje 'no detectado' NO implica que no ocurrió — solo que no dejó huella en FISSAL.", 0.5, 5.5, 12, 0.9)

# =====================================================================
# SLIDE 4: HITO 2 — DIAGNÓSTICO
# =====================================================================
slide = add_slide()
add_title(slide, "Hito 2: Diagnóstico")
add_subtitle(slide, "Primera atención codificada como cáncer colorrectal (C18/C19/C20)")
add_footer(slide)

# Tipo de atención al diagnóstico
tipos_dx = hitos["TIPO_ATENCION_DIAGNOSTICO"].value_counts()
rows_tipos = [[t, n, f"{n/N_PAC*100:.1f}%"] for t, n in tipos_dx.items()]
add_text(slide, "Tipo de atención en el momento del diagnóstico", 0.5, 1.6, 6, 0.4, font_size=13, bold=True)
add_table(slide, ["Tipo", "Pacientes", "%"], rows_tipos, 0.5, 2.1, [2.5, 1.5, 1.0])

# CIE-10 prefix
hitos["CIE10_PREFIJO"] = hitos["CIE10_DIAGNOSTICO"].str[:3]
cie10 = hitos["CIE10_PREFIJO"].value_counts()
rows_cie = [[c, n, f"{n/N_PAC*100:.1f}%"] for c, n in cie10.items()]
add_text(slide, "Localización anatómica (prefijo CIE-10)", 6.8, 1.6, 5, 0.4, font_size=13, bold=True)
add_table(slide, ["CIE-10", "Pacientes", "%"], rows_cie, 6.8, 2.1, [1.5, 1.5, 1.0])

# Diagnóstico por año
hitos["ANIO_DX"] = pd.to_datetime(hitos["FECHA_DIAGNOSTICO"]).dt.year
anios_dx = hitos["ANIO_DX"].value_counts().sort_index()
rows_yr = [[int(yr), n] for yr, n in anios_dx.items()]
add_text(slide, "Diagnósticos por año (fecha real de atención)", 0.5, 3.7, 6, 0.4, font_size=13, bold=True)
add_table(slide, ["Año", "Pacientes"], rows_yr, 0.5, 4.2, [1.5, 1.5])

# Urgencia por edad
add_text(slide, "Diagnóstico urgente (emergencia/hosp) por edad", 6.8, 3.7, 5, 0.4, font_size=13, bold=True)
hitos["INGRESO_URGENTE"] = hitos["TIPO_ATENCION_DIAGNOSTICO"].isin(["EMERGENCIA", "HOSPITALIZACIÓN"])
rangos = ["0-17", "18-29", "30-39", "40-49", "50-59", "60-69", "70-79", "80+"]
rows_urg = []
for r in rangos:
    sub = hitos[hitos["RANGO_EDAD"] == r]
    if len(sub) == 0: continue
    n_urg = sub["INGRESO_URGENTE"].sum()
    rows_urg.append([r, len(sub), n_urg, f"{n_urg/len(sub)*100:.1f}%"])
add_table(slide, ["Edad", "n", "Urgente", "%"], rows_urg, 6.8, 4.2, [1.0, 1.0, 1.0, 1.0])

add_box(slide, "Diagnóstico vía EMERGENCIA u HOSPITALIZACIÓN sugiere un cuadro más urgente/avanzado. Diagnóstico vía consulta AMBULATORIA sugiere detección menos urgente. Ambos tipos son igualmente válidos como fecha de diagnóstico formal.", 0.5, 6.3, 12, 0.6)

# =====================================================================
# SLIDE 5: HITO 3 — TRATAMIENTO (modalidades)
# =====================================================================
slide = add_slide()
add_title(slide, "Hito 3: Tratamiento — Modalidades")
add_subtitle(slide, "Cirugía, Quimioterapia y Radioterapia detectadas por patrones de texto en descripciones")
add_footer(slide)

# KPIs de modalidades
for i, cat in enumerate(["CIRUGIA", "QUIMIOTERAPIA", "RADIOTERAPIA"]):
    col = f"TUVO_{cat}"
    n = hitos[col].sum()
    add_kpi(slide, cat, f"{n:,} ({n/N_PAC*100:.1f}%)", 0.5 + i * 4.15, 1.7, width=3.8, color=[PRIMARY, MID_BLUE, NEUTRAL][i])

# Combinaciones
add_text(slide, "Secuencias de tratamiento más frecuentes", 0.5, 3.5, 6, 0.4, font_size=13, bold=True)
combos = hitos["SECUENCIA_TRATAMIENTO"].value_counts().head(10)
rows_combos = [[c, n, f"{n/N_PAC*100:.1f}%"] for c, n in combos.items()]
add_table(slide, ["Secuencia", "Pacientes", "%"], rows_combos, 0.5, 4.0, [4.0, 1.5, 1.0])

add_box(slide, "El 68.1% de pacientes NO tiene tratamiento detectado por estos patrones. Posibles razones: (a) códigos sin descripción reconocible, (b) pacientes con solo manejo sintomático/paliativo sin cirugía ni quimio, (c) tratamientos realizados fuera de FISSAL.", 6.8, 3.5, 6.0, 1.2)

# Tiempo diagnóstico → tratamiento
tt = hitos["DIAS_DIAGNOSTICO_A_TRATAMIENTO"].dropna()
tt = tt[tt >= 0]
add_kpi(slide, "Mediana dx → tratamiento", f"{tt.median():.0f} días", 0.5, 5.6, color=PRIMARY)
add_kpi(slide, "Media dx → tratamiento", f"{tt.mean():.0f} días", 3.6, 5.6, color=MID_BLUE)
add_text(slide, "* Tiempo de oportunidad: mide cuánto tarda el sistema en iniciar tratamiento desde el diagnóstico formal", 0.5, 6.4, 8, 0.3, font_size=8, color=GRAY)

# =====================================================================
# SLIDE 5b: ¿QUÉ PASA CON EL 68.1% SIN TRATAMIENTO DETECTADO?
# =====================================================================
slide = add_slide()
add_title(slide, "¿Qué pasa con el 68.1% sin tratamiento detectado?")
add_subtitle(slide, "Análisis detallado de los 10,691 pacientes que NO tienen cirugía, quimioterapia ni radioterapia identificada")
add_footer(slide)

sin = hitos[hitos["SECUENCIA_TRATAMIENTO"] == "SIN_TRATAMIENTO_DETECTADO"]
con = hitos[hitos["SECUENCIA_TRATAMIENTO"] != "SIN_TRATAMIENTO_DETECTADO"]

add_kpi(slide, "Sin tratamiento (68.1%)", f"{len(sin):,}", 0.5, 1.6, width=3.8, color=BAD)
add_kpi(slide, "Con tratamiento (31.9%)", f"{len(con):,}", 6.7, 1.6, width=3.8, color=GOOD)

# Qué SÍ tienen
add_text(slide, "Estos 10,691 pacientes SÍ tienen actividad en FISSAL — simplemente no coincide con los patrones de cirugía/quimio/radio:", 0.5, 3.3, 12, 0.4, font_size=11)
add_kpi(slide, "Procedimientos\n(pero no quirúrgicos)", f"{9869} (92.3%)", 0.5, 3.9, width=2.8, height=1.2)
add_kpi(slide, "Medicamentos\n(pero no quimio)", f"{6266} (58.6%)", 3.6, 3.9, width=2.8, height=1.2)
add_kpi(slide, "Hospitalización", f"{1446} (13.5%)", 6.7, 3.9, width=2.8, height=1.2)
add_kpi(slide, "Emergencia", f"{1100} (10.3%)", 9.8, 3.9, width=2.8, height=1.2)

# Comparativa
add_text(slide, "Comparación de intensidad de atención: SIN vs CON tratamiento", 0.5, 5.4, 12, 0.4, font_size=12, bold=True)
rows_comp = [
    ["Nº atenciones (mediana)", "1", "15", "Los SIN tratamiento típicamente tienen 1 sola visita (la del diagnóstico)"],
    ["Nº prestaciones (mediana)", "7", "184", "Solo exámenes/labs básicos vs. seguimiento oncológico completo"],
    ["Tiempo en sistema (mediana)", "0 días", "321 días", "El paciente SIN tratamiento no vuelve después del diagnóstico"],
    ["Costo neto total (mediana)", "S/ 102", "S/ 4,457", "Costo 44x menor — coherente con atención mínima"],
]
add_table(slide, ["Métrica", "SIN tratamiento", "CON tratamiento", "Interpretación"], rows_comp, 0.5, 5.9, [2.5, 1.5, 1.5, 5.5], font_size=9)

# =====================================================================
# SLIDE 5c: DESENLACE DE LOS SIN TRATAMIENTO
# =====================================================================
slide = add_slide()
add_title(slide, "Desenlace de los Pacientes sin Tratamiento Detectado")
add_subtitle(slide, "¿Qué pasó con esos 10,691 pacientes después del diagnóstico?")
add_footer(slide)

# KPIs de desenlace sin tratamiento
add_kpi(slide, "Posible remisión/alta", "6,028 (56.4%)", 0.5, 1.6, color=GOOD)
add_kpi(slide, "Seguimiento activo", "2,788 (26.1%)", 3.6, 1.6, color=PRIMARY)
add_kpi(slide, "Fallecidos", "1,461 (13.7%)", 6.7, 1.6, color=BAD)
add_kpi(slide, "Control post-pausa", "387 (3.6%)", 9.8, 1.6, color=MID_BLUE)

# Explicacion por categoria
add_text(slide, "¿Qué significa cada desenlace en este grupo?", 0.5, 3.3, 12, 0.4, font_size=13, bold=True)

explicaciones_sin = [
    ("POSIBLE_REMISIÓN_O_ALTA (56.4%)", "La gran mayoría. Tuvieron una o pocas visitas y luego desaparecieron de FISSAL por ≥18 meses. Pueden haber sido derivados a otro sistema (EsSalud, MINSA), haber fallecido sin registro, o tener enfermedad muy temprana que no requirió tratamiento oncológico en FISSAL."),
    ("SEGUIMIENTO ACTIVO (26.1%)", "Siguen apareciendo en FISSAL sin brechas largas. Van a consultas, laboratorios, imágenes... pero ningún registro coincide con los patrones de cirugía/quimio/radio. Podrían estar en vigilancia activa post-polipectomía o con enfermedad indolente."),
    ("FALLECIDOS (13.7%)", "Murieron durante el seguimiento. De ellos, la mayoría probablemente falleció poco después del diagnóstico sin llegar a recibir tratamiento oncológico (enfermedad muy avanzada al debut, o fallecieron por otra causa)."),
    ("CONTROL POST-PAUSA (3.6%)", "Se fueron y volvieron solo para chequeo. Compatible con pacientes derivados que regresan a FISSAL para control de rutina."),
    ("RECAIDA PROBABLE (0.2%)", "Solo 17 pacientes. Se fueron ≥18 meses y volvieron con tratamiento activo. Estos SÍ recibieron tratamiento, pero en otro sistema — FISSAL solo ve su regreso."),
]

y = 3.9
for tit, desc in explicaciones_sin:
    add_bullet(slide, f"▪ {tit}", 0.5, y, 12, 0.3, font_size=11, bold=True, color=PRIMARY)
    add_text(slide, desc, 0.8, y + 0.3, 11.7, 0.45, font_size=9, color=GRAY)
    y += 0.75

add_box(slide, "Conclusión: 'SIN_TRATAMIENTO_DETECTADO' NO significa que el paciente no recibió tratamiento. Significa que FISSAL no tiene registro de cirugía/quimio/radio CON descripciones que coincidan con nuestros patrones de búsqueda. El 56% simplemente desapareció del sistema después del diagnóstico — probablemente derivados a otra institución.", 0.5, 6.3, 12, 0.6)

# =====================================================================
# SLIDE 6: HITO 3 — ADHERENCIA Y HOSPITALIZACIÓN
# =====================================================================
slide = add_slide()
add_title(slide, "Hito 3: Adherencia y Hospitalización durante el Tratamiento")
add_subtitle(slide, "Brechas de atención (>45 días) y episodios de hospitalización")
add_footer(slide)

con_trat = hitos[hitos["SECUENCIA_TRATAMIENTO"] != "SIN_TRATAMIENTO_DETECTADO"]
n_brecha = con_trat["TUVO_BRECHA_ADHERENCIA"].sum()
n_hosp = (hitos["N_HOSPITALIZACIONES"] > 0).sum()

add_kpi(slide, "Pacientes con ≥1 brecha (>45d)", f"{n_brecha:,} ({n_brecha/len(con_trat)*100:.1f}%)", 0.5, 1.7, color=NEUTRAL)
add_kpi(slide, "Pacientes hospitalizados", f"{n_hosp:,} ({n_hosp/N_PAC*100:.1f}%)", 3.6, 1.7, color=MID_BLUE)

# Adherencia por edad
add_text(slide, "% con brechas de adherencia por rango de edad", 0.5, 3.4, 5, 0.4, font_size=13, bold=True)
rows_brecha = []
for r in rangos:
    sub = con_trat[con_trat["RANGO_EDAD"] == r]
    if len(sub) < 5: continue
    nb = sub["TUVO_BRECHA_ADHERENCIA"].sum()
    rows_brecha.append([r, len(sub), nb, f"{nb/len(sub)*100:.1f}%"])
add_table(slide, ["Edad", "n", "Con brecha", "%"], rows_brecha, 0.5, 3.9, [1.0, 1.2, 1.3, 1.0])

# Hospitalización por desenlace
add_text(slide, "% hospitalizados por tipo de desenlace", 6.8, 3.4, 5, 0.4, font_size=13, bold=True)
rows_hosp_d = []
for cat in ["FALLECIDO", "PALIATIVO", "RECAIDA_PROBABLE", "CONTROL_POST_PAUSA", "POSIBLE_REMISION_O_ALTA", "EN_SEGUIMIENTO_ACTIVO"]:
    sub = hitos[hitos["DESENLACE"] == cat]
    if len(sub) == 0: continue
    n_h = (sub["N_HOSPITALIZACIONES"] > 0).sum()
    rows_hosp_d.append([cat.replace("_", " "), len(sub), n_h, f"{n_h/len(sub)*100:.1f}%"])
add_table(slide, ["Desenlace", "n", "Hosp.", "%"], rows_hosp_d, 6.8, 3.9, [3.0, 1.0, 1.0, 1.0])

add_box(slide, "Brecha de adherencia = vacío >45 días entre atenciones consecutivas DENTRO del tramo de tratamiento activo. Los esquemas de quimio típicos (FOLFOX/FOLFIRI) ciclan cada 14-21 días; >45d duplica el ciclo más largo esperado. FISSAL no registra citas programadas vs. perdidas — es una aproximación.", 0.5, 5.9, 12, 0.8)

# =====================================================================
# SLIDE 7: HITO 4 — DESENLACE
# =====================================================================
slide = add_slide()
add_title(slide, "Hito 4: Desenlace — Clasificación")
add_subtitle(slide, "6 categorías inferidas a partir de fallecimiento, cuidados paliativos, y brechas de atención")
add_footer(slide)

desenlaces_data = [
    ("FALLECIDO", BAD, "Fecha de fallecimiento registrada en FISSAL"),
    ("PALIATIVO", NEUTRAL, "Recibió visita domiciliaria de cuidados paliativos (2016-2018)"),
    ("RECAIDA_PROBABLE", RGBColor(0xD7, 0x30, 0x27), "Brecha ≥540d y volvió con cirugía/quimio/radio/hosp"),
    ("CONTROL_POST_PAUSA", RGBColor(0x45, 0x75, 0xB4), "Brecha ≥540d y volvió solo con consulta/lab/imagen"),
    ("POSIBLE_REMISION_O_ALTA", GOOD, "Brecha ≥540d es lo ÚLTIMO que se ve — nunca volvió"),
    ("EN_SEGUIMIENTO_ACTIVO", PRIMARY, "Sin brechas ≥540d en toda su historia"),
]

desenlace_counts = hitos["DESENLACE"].value_counts()
rows_des = []
for cat, color, desc in desenlaces_data:
    n = desenlace_counts.get(cat, 0)
    rows_des.append([cat.replace("_", " "), n, f"{n/N_PAC*100:.1f}%", desc])

add_text(slide, "Distribución de desenlaces", 0.5, 1.5, 6, 0.4, font_size=14, bold=True)
add_table(slide, ["Desenlace", "Pacientes", "%", "Significado"], rows_des, 0.5, 2.0, [2.5, 1.3, 1.0, 5.5], font_size=9)

# KPIs
add_kpi(slide, "Fallecidos", f"{desenlace_counts.get('FALLECIDO',0):,}", 0.5, 4.7, color=BAD)
add_kpi(slide, "Posible remisión/alta", f"{desenlace_counts.get('POSIBLE_REMISION_O_ALTA',0):,}", 3.6, 4.7, color=GOOD)
add_kpi(slide, "Recaída probable", f"{desenlace_counts.get('RECAIDA_PROBABLE',0):,}", 6.7, 4.7, color=RGBColor(0xD7, 0x30, 0x27))
add_kpi(slide, "Seguimiento activo", f"{desenlace_counts.get('EN_SEGUIMIENTO_ACTIVO',0):,}", 9.8, 4.7, color=PRIMARY)

add_box(slide, "FISSAL no tiene campo de 'alta' o 'recuperación'. POSIBLE_REMISION_O_ALTA es una inferencia por inactividad (pausa ≥540d sin retorno). Puede ser remisión real, cambio de aseguradora, migración o pérdida de seguimiento. NO es confirmación clínica.", 0.5, 6.3, 12, 0.6)

# =====================================================================
# SLIDE 8: DESENLACE — CARACTERÍSTICAS POR GRUPO
# =====================================================================
slide = add_slide()
add_title(slide, "Desenlace: Características por Grupo")
add_subtitle(slide, "Perfil demográfico, tratamiento y costos de cada categoría de desenlace")
add_footer(slide)

rows_perfil = []
for cat, color, _ in desenlaces_data:
    sub = hitos[hitos["DESENLACE"] == cat]
    if len(sub) == 0: continue
    edad = sub["EDAD_PRIMERA_ATENCION"].dropna()
    pct_f = sub[sub["SEXO"] == "F"].shape[0] / len(sub) * 100
    pct_cir = sub["TUVO_CIRUGIA"].mean() * 100
    pct_qt = sub["TUVO_QUIMIOTERAPIA"].mean() * 100
    pct_hosp = (sub["N_HOSPITALIZACIONES"] > 0).mean() * 100
    costo_med = sub["MONTO_NETO_TOTAL"].median()
    rows_perfil.append([
        cat.replace("_", " "), len(sub), f"{edad.mean():.0f}" if len(edad)>0 else "N/D",
        f"{pct_f:.0f}%", f"{pct_cir:.0f}%", f"{pct_qt:.0f}%", f"{pct_hosp:.0f}%", f"S/{costo_med:,.0f}"
    ])

add_text(slide, "Perfil comparativo por tipo de desenlace", 0.5, 1.5, 6, 0.4, font_size=13, bold=True)
add_table(slide, ["Desenlace", "n", "Edad media", "%F", "%Cir", "%QT", "%Hosp", "Costo med."], rows_perfil,
          0.5, 2.0, [2.5, 1.0, 1.0, 0.7, 0.7, 0.7, 0.8, 1.5], font_size=9)

# Hallazgos clave
add_text(slide, "Hallazgos clave", 0.5, 5.0, 6, 0.3, font_size=13, bold=True)
add_bullet(slide, "RECAIDA_PROBABLE tiene la tasa de hospitalización más alta de todos los grupos (65.7%) — su regreso implicó atención más intensiva.", 0.5, 5.4, 12, 0.35, font_size=10)
add_bullet(slide, "POSIBLE_REMISION_O_ALTA es el grupo más numeroso (50.9%) y el de menor costo mediano (S/214), coherente con pacientes que dejaron de necesitar atención.", 0.5, 5.75, 12, 0.35, font_size=10)
add_bullet(slide, "PALIATIVO solo tiene 19 pacientes después de la corrección de outliers en costo; antes de corregir, el costo medio estaba inflado por un outlier de S/723M.", 0.5, 6.1, 12, 0.35, font_size=10)
add_bullet(slide, "La gran mayoría de pausas largas (≥540d) son el FINAL de la historia del paciente. Solo 762 pacientes (4.8%) 'se fueron y volvieron' — los RECAIDA_PROBABLE + CONTROL_POST_PAUSA.", 0.5, 6.45, 12, 0.35, font_size=10)

# =====================================================================
# SLIDE 9: TIEMPOS TRANSVERSALES
# =====================================================================
slide = add_slide()
add_title(slide, "Tiempos en la Trayectoria Completa")
add_subtitle(slide, "Duración de cada tramo entre hitos, comparado por sexo y edad")
add_footer(slide)

# Calcular tiempos
intervalos = [
    ("DIAS_DESPISTAJE_A_DIAGNOSTICO", "Despistaje → Diagnóstico", "Solo con despistaje detectado"),
    ("DIAS_DIAGNOSTICO_A_TRATAMIENTO", "Diagnóstico → Inicio Tratamiento", "Solo con tratamiento detectado"),
    ("DURACION_CIRUGIA_DIAS", "Duración Cirugía (1ra-última sesión)", "Solo con cirugía"),
    ("DURACION_QUIMIOTERAPIA_DIAS", "Duración Quimioterapia", "Solo con quimioterapia"),
]

rows_int = []
for col, tit, nota in intervalos:
    s = hitos[col].dropna()
    s = s[s >= 0]
    if len(s) == 0: continue
    rows_int.append([tit, len(s), f"{s.median():.0f}", f"{s.mean():.0f}", f"{s.quantile(.25):.0f}", f"{s.quantile(.75):.0f}", nota])

add_text(slide, "Estadísticos de duración por tramo (días)", 0.5, 1.5, 6, 0.4, font_size=13, bold=True)
add_table(slide, ["Tramo", "n", "Mediana", "Media", "P25", "P75", "Nota"], rows_int,
          0.5, 2.0, [3.0, 1.0, 1.0, 1.0, 1.0, 1.0, 3.0], font_size=9)

# Tiempo dx→tratamiento por sexo
add_text(slide, "Tiempo Diagnóstico → Tratamiento por sexo", 0.5, 4.2, 6, 0.4, font_size=13, bold=True)
rows_sex = []
for sexo in ["F", "M"]:
    s = hitos.loc[hitos["SEXO"] == sexo, "DIAS_DIAGNOSTICO_A_TRATAMIENTO"].dropna()
    s = s[s >= 0]
    rows_sex.append([sexo, len(s), f"{s.median():.0f}", f"{s.mean():.0f}"])
add_table(slide, ["Sexo", "n", "Mediana (d)", "Media (d)"], rows_sex, 0.5, 4.7, [1.5, 1.5, 1.5, 1.5])

# Tiempo dx→tratamiento por edad
add_text(slide, "Tiempo Diagnóstico → Tratamiento por rango de edad", 6.8, 4.2, 5, 0.4, font_size=13, bold=True)
rows_edad_t = []
for r in rangos:
    s = hitos.loc[hitos["RANGO_EDAD"] == r, "DIAS_DIAGNOSTICO_A_TRATAMIENTO"].dropna()
    s = s[s >= 0]
    if len(s) < 5: continue
    rows_edad_t.append([r, len(s), f"{s.median():.0f}", f"{s.mean():.0f}"])
add_table(slide, ["Edad", "n", "Mediana", "Media"], rows_edad_t, 6.8, 4.7, [1.0, 1.0, 1.0, 1.0])

add_box(slide, "El tiempo diagnóstico → tratamiento es el indicador clave de OPORTUNIDAD de atención. Una mediana baja sugiere que el sistema responde rápido; una mediana alta puede indicar barreras de acceso, listas de espera o derivaciones tardías.", 0.5, 6.2, 12, 0.7)

# =====================================================================
# SLIDE 10: CALIDAD DE DATOS
# =====================================================================
slide = add_slide()
add_title(slide, "Calidad de Datos: Correcciones Aplicadas")
add_subtitle(slide, "Outliers de cantidad, costo y edad que fueron detectados y corregidos")
add_footer(slide)

add_text(slide, "1. Outliers de cantidad (y por tanto de costo)", 0.5, 1.5, 6, 0.4, font_size=14, bold=True, color=BAD)
add_text(slide, "415 registros corregidos por cantidades absurdas (ej. 1,112,500 grapadoras quirúrgicas). Se usó la mediana histórica por código de ítem como referencia. Impacto: S/861.5 millones en montos absurdos reescalados a valores razonables. El caso más extremo: filtro de ventilador facturado en S/80M (3.3M unidades, mediana del ítem = 1).", 0.5, 2.0, 12, 1.0, font_size=10)

add_text(slide, "2. Edades clínicamente implausibles para CCR", 0.5, 3.3, 6, 0.4, font_size=14, bold=True, color=BAD)
add_text(slide, "75 pacientes con edad <18 años (incluyendo 7 con edad 0-3 años). El cáncer colorrectal en primera infancia es prácticamente inexistente. Se marcaron como NA las edades <10 años (umbral conservador) y >100 años. Resultado: 19 pacientes con edad inválida excluidos de los análisis de edad.", 0.5, 3.8, 12, 0.8, font_size=10)

add_text(slide, "3. Costo: preferir siempre la MEDIANA sobre la media", 0.5, 4.9, 6, 0.4, font_size=14, bold=True, color=PRIMARY)
add_text(slide, "Ejemplo: el costo neto medio de POSIBLE_REMISION_O_ALTA bajó de S/90,967 a S/3,444 después de corregir outliers. La mediana casi no se movió (S/220 → S/214), confirmando que ya era confiable antes de la corrección.", 0.5, 5.4, 12, 0.6, font_size=10)

add_text(slide, "4. Oxígeno medicinal: la categoría menos certera", 0.5, 6.1, 6, 0.3, font_size=13, bold=True, color=NEUTRAL)
add_text(slide, "253 de los 415 registros corregidos son de oxígeno medicinal. Su precio unitario varía mucho (S/0.001 a S/10 por 'unidad'), lo que sugiere que el ítem mezcla distintas unidades de medida bajo el mismo código. No hay certeza absoluta en estas correcciones.", 0.5, 6.4, 12, 0.5, font_size=9, color=GRAY)

# =====================================================================
# SLIDE 11: LIMITACIONES
# =====================================================================
slide = add_slide()
add_title(slide, "Limitaciones Metodológicas")
add_subtitle(slide, "Consideraciones importantes para interpretar los resultados")
add_footer(slide)

limitaciones = [
    ("Despistaje invisible fuera de FISSAL", "El 93.4% de pacientes aparece en FISSAL ya con diagnóstico. Esto NO significa ausencia de despistaje: probablemente ocurrió en EsSalud, consulta privada o MINSA sin cobertura FISSAL."),
    ("Tratamiento no detectado (68.1%)", "Los patrones de cirugía/quimio/radio son aproximaciones por texto. Puede haber tratamientos con códigos sin descripción reconocible, o pacientes que solo reciben manejo sintomático."),
    ("POSIBLE_REMISION_O_ALTA es una hipótesis", "Inferida por inactividad ≥540 días. Puede ser remisión real, cambio de aseguradora, migración, abandono o muerte no registrada. Tratarla siempre como hipótesis, no como hecho."),
    ("Adherencia mide vacíos en facturación", "FISSAL no tiene agenda de citas. Un vacío >45d en la facturación no necesariamente es una inasistencia. Puede ser que el paciente recibió atención pero no se facturó, o que el esquema de tratamiento lo permitía."),
    ("Recaída probable no es confirmación clínica", "Se infiere cuando un paciente con brecha ≥540d regresa con tratamiento activo. Pudiera ser un nuevo primario, progresión de enfermedad preexistente o complicación tardía."),
    ("Edades <10 años marcadas como NA", "Filtro conservador específico para CCR. Puede descartar algún caso pediátrico real (poliposis adenomatosa familiar), pero la probabilidad es ínfima."),
]

y = 1.6
for tit, desc in limitaciones:
    add_bullet(slide, f"▪ {tit}", 0.7, y, 11.5, 0.35, font_size=12, bold=True, color=BAD)
    add_text(slide, desc, 1.1, y + 0.3, 11.3, 0.45, font_size=9, color=GRAY)
    y += 0.85

# =====================================================================
# GUARDAR
# =====================================================================
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(PPTX_OUT))
print(f"\nPresentacion guardada: {PPTX_OUT}")
print(f"Total slides: {len(prs.slides)}")

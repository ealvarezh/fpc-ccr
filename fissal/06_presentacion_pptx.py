import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import gc

# =====================================================================
# CONFIGURACION
# =====================================================================
SILVER = Path(r"C:\Users\eah\apoyoconsultoria.com\File Server - Analytics\7 Datos\Datos abiertos\fissal\01_silver")
OUTPUT_DIR = SILVER.parent / "03_output" / "supervivencia"
PPTX_OUT = OUTPUT_DIR.parent / "FISSAL_CCR_Presentacion.pptx"

PRIMARY   = RGBColor(0x11, 0x8D, 0xFF)
FOREGROUND= RGBColor(0x25, 0x24, 0x23)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF3, 0xF2, 0xF1)
GOOD      = RGBColor(0x1A, 0xAB, 0x40)
NEUTRAL   = RGBColor(0xD9, 0xB3, 0x00)
BAD       = RGBColor(0xD6, 0x45, 0x54)
DARK_BLUE = RGBColor(0x0D, 0x5E, 0xAA)
MID_BLUE  = RGBColor(0x43, 0x93, 0xC3)
GRAY      = RGBColor(0x60, 0x5E, 0x5C)
LIGHTER   = RGBColor(0xC0, 0xD0, 0xE8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# =====================================================================
# CARGA DE DATOS
# =====================================================================
print("Cargando datos...")
perfil = pd.read_parquet(SILVER / "FISSAL_CCR_PERFIL_PACIENTES.parquet")
perfil_amp = pd.read_parquet(SILVER / "FISSAL_CCR_PERFIL_PACIENTES_AMPLIADO.parquet")

# ===== FILTRAR OUTLIERS DE EDAD =====
perfil.loc[(perfil["EDAD_PRIMERA_ATENCION"] < 0) | (perfil["EDAD_PRIMERA_ATENCION"] > 120), "EDAD_PRIMERA_ATENCION"] = pd.NA

# Cargar datos completos CCR para trayectoria clinica
print("Cargando datos CCR para trayectoria clinica...")
df_ccr = pd.read_parquet(SILVER / "FISSAL_CANCER_COLORRECTAL_2016_2022.parquet",
                         columns=["NUM_DOC_CRYPTO", "ATE_FECATENCION", "ATE_MONTONETO", "ATE_MONTOBRUTO",
                                  "ATE_TIPOCONSUMO", "ATE_DESCCONSUMO", "TIPO_ATENC",
                                  "ATE_FECINGHOSP", "ATE_FECALTHOSP", "ANIO_ATENCION"])

N_PAC = len(perfil)
N_FALL = perfil["FALLECIDO"].sum()
N_VIVOS = N_PAC - N_FALL
PCT_FALL = N_FALL / N_PAC * 100

# =====================================================================
# FUNCIONES AUXILIARES
# =====================================================================
def add_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def add_title(slide, text, left=0.5, top=0.3, width=12.3, height=0.7, font_size=28):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = True; p.font.color.rgb = FOREGROUND; p.font.name = "Segoe UI"
    return txBox

def add_subtitle(slide, text, left=0.5, top=1.0, width=12.3, height=0.4, font_size=16):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.color.rgb = GRAY; p.font.name = "Segoe UI"
    return txBox

def add_text_box(slide, text, left, top, width, height, font_size=12, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    if color is None: color = FOREGROUND
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = "Segoe UI"
    p.alignment = alignment
    return txBox

def add_bullet(slide, text, left, top, width, height, font_size=11, color=None, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color or FOREGROUND; p.font.name = "Segoe UI"
    return txBox

def add_kpi_card(slide, label, value, left, top, width=2.8, height=1.4, color=None):
    if color is None: color = PRIMARY
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE1, 0xDF, 0xDD); shape.line.width = Pt(1)
    add_text_box(slide, str(value), left + 0.2, top + 0.15, width - 0.4, 0.65, font_size=28, bold=True, color=color)
    add_text_box(slide, label, left + 0.2, top + 0.8, width - 0.4, 0.4, font_size=11, color=GRAY)

def add_table(slide, headers, rows, left, top, col_widths, font_size=10):
    n_rows = len(rows) + 1; n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(total_w), Inches(0.35 * n_rows))
    table = table_shape.table
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
        for p in table.cell(0, j).text_frame.paragraphs:
            p.font.size = Pt(font_size); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
        table.cell(0, j).fill.solid(); table.cell(0, j).fill.fore_color.rgb = PRIMARY
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            table.cell(i + 1, j).text = str(val)
            for p in table.cell(i + 1, j).text_frame.paragraphs:
                p.font.size = Pt(font_size); p.font.color.rgb = FOREGROUND; p.font.name = "Segoe UI"
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            table.cell(i + 1, j).fill.solid()
            table.cell(i + 1, j).fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_BG
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)
    return table_shape

def add_footer(slide, text="FISSAL CCR | Atenciones 2012-2022 | Datos procesados 2016-2022"):
    add_text_box(slide, text, 0.5, 7.0, 12, 0.3, font_size=8, color=RGBColor(0xA1, 0x9F, 0x9D))

def add_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = PRIMARY; shape.line.fill.background()

def add_explanation_box(slide, text, left, top, width, height=0.6):
    """Caja de explicacion con fondo gris claro."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    add_text_box(slide, text, left + 0.15, top + 0.08, width - 0.3, height - 0.15, font_size=9, color=GRAY)

# =====================================================================
# SLIDE 1: PORTADA
# =====================================================================
print("Slides...", end=" ")
slide = add_slide(prs)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE
add_text_box(slide, "FISSAL  |  CÁNCER COLORRECTAL", 1, 1.8, 11, 1.2, font_size=38, bold=True, color=WHITE)
add_text_box(slide, "Supervivencia, Tratamiento y Trayectoria Clínica", 1, 3.0, 11, 0.8, font_size=20, color=LIGHTER)
add_line(slide, 1, 3.9, 4)
add_text_box(slide, "Atenciones 2012-2022  |  15,698 pacientes  |  CIE-10: C18/C19/C20  |  FISSAL - Perú", 1, 4.3, 11, 0.5, font_size=14, color=LIGHTER)
add_text_box(slide, "Fuente: Fondo Intangible Solidario de Salud — Prestaciones procesadas 2016-2022 (fechas reales de atención: 2012-2022)", 1, 5.2, 11, 0.5, font_size=11, color=LIGHTER)

# =====================================================================
# SLIDE 2: RESUMEN GENERAL
# =====================================================================
slide = add_slide(prs)
add_title(slide, "Resumen de la Cohorte CCR")
add_subtitle(slide, "Pacientes con cáncer colorrectal (CIE-10: C18, C19, C20). Atenciones reales: 2012-2022. Datos procesados por FISSAL: 2016-2022.")
add_footer(slide)

add_kpi_card(slide, "Pacientes totales", f"{N_PAC:,}", 0.5, 1.8, color=PRIMARY)
add_kpi_card(slide, "Prestaciones totales", "1,559,428", 3.6, 1.8, color=MID_BLUE)
add_kpi_card(slide, "Fallecidos registrados", f"{N_FALL:,} ({PCT_FALL:.1f}%)", 6.7, 1.8, color=BAD)
add_kpi_card(slide, "Mediana de edad", f"{perfil['EDAD_PRIMERA_ATENCION'].median():.0f} años", 9.8, 1.8, color=GOOD)

anios_data = perfil["PRIMER_ANIO"].value_counts().sort_index()
rows_anio = []
for yr in anios_data.index:
    yr_int = int(yr)
    n_nuevos = anios_data[yr]
    n_total = perfil[perfil["ANIOS"].astype(str).str.contains(str(yr_int))].shape[0]
    rows_anio.append([yr_int, n_nuevos, n_total, f"{n_nuevos/n_total*100:.1f}%"])

add_text_box(slide, "Pacientes por año de primera atención (fecha real, no producción)", 0.5, 3.5, 5, 0.4, font_size=12, bold=True)
add_table(slide, ["Año", "Nuevos", "Total activos", "% Nuevos"], rows_anio, 0.5, 4.0, [0.8, 1.2, 1.2, 0.9], font_size=9)

sexo = perfil["SEXO"].value_counts()
rows_sexo = [[s, n, f"{n/N_PAC*100:.1f}%"] for s, n in sexo.items()]
add_text_box(slide, "Distribución por sexo", 6.8, 3.6, 5, 0.4, font_size=13, bold=True)
add_table(slide, ["Sexo", "Pacientes", "%"], rows_sexo, 6.8, 4.1, [1.5, 1.5, 1.0], font_size=10)

# Explicacion
add_explanation_box(slide, "Los años 2012-2015 provienen de registros históricos cargados en FISSAL. La producción (procesamiento) de datos por FISSAL empezó en 2016. La fecha que importa es la de ATENCIÓN real (ATE_FECATENCION), no la de producción (ATE_ANOPROD).", 0.5, 5.8, 12, 0.7)

# =====================================================================
# SLIDE 3: DEMOGRAFIA
# =====================================================================
slide = add_slide(prs)
add_title(slide, "Perfil Demográfico")
add_subtitle(slide, "Edad a la primera atención (con filtro de outliers: 0 a 120 años)")
add_footer(slide)

edad = perfil["EDAD_PRIMERA_ATENCION"].dropna()
add_kpi_card(slide, "Edad media", f"{edad.mean():.1f} años", 0.5, 1.8, color=PRIMARY)
add_kpi_card(slide, "Edad mediana", f"{edad.median():.0f} años", 3.6, 1.8, color=MID_BLUE)
add_kpi_card(slide, "Mínima", f"{edad.min():.0f} años", 6.7, 1.8, color=GOOD)
add_kpi_card(slide, "Máxima", f"{edad.max():.0f} años", 9.8, 1.8, color=BAD)

n_outliers = perfil["EDAD_PRIMERA_ATENCION"].isna().sum()
add_text_box(slide, f"Registros con edad inválida (negativa o >120) excluidos: {n_outliers}", 0.5, 3.1, 6, 0.3, font_size=9, color=BAD)

bins = [0, 18, 30, 40, 50, 60, 70, 80, 120]
labels = ["0-17", "18-29", "30-39", "40-49", "50-59", "60-69", "70-79", "80+"]
temp_rango = pd.cut(edad, bins=bins, labels=labels, right=False)
rango_counts = temp_rango.value_counts().sort_index()

rows_rango = []
for r in labels:
    n = rango_counts.get(r, 0)
    rows_rango.append([r, n, f"{n/len(edad)*100:.1f}%"])

add_text_box(slide, "Distribución por rango de edad", 0.5, 3.6, 5, 0.4, font_size=13, bold=True)
add_table(slide, ["Rango", "Pacientes", "%"], rows_rango, 0.5, 4.1, [1.5, 1.5, 1.0], font_size=10)

dep = perfil["DEPARTAMENTO"].value_counts().head(10)
rows_dep = [[d, n, f"{n/N_PAC*100:.1f}%"] for d, n in dep.items()]
add_text_box(slide, "Top 10 departamentos de residencia", 6.8, 3.6, 5, 0.4, font_size=13, bold=True)
add_table(slide, ["Departamento", "Pacientes", "%"], rows_dep, 6.8, 4.1, [2.5, 1.5, 1.0], font_size=10)

add_explanation_box(slide, "La edad se calcula como floor((fecha 1ra atención - fecha nacimiento) / 365.25). Se excluyen edades negativas (fecha de nacimiento posterior a la atención) y >120 años (errores de registro).", 0.5, 6.2, 12, 0.5)

# =====================================================================
# SLIDE 4: MORTALIDAD
# =====================================================================
slide = add_slide(prs)
add_title(slide, "Mortalidad: Tiempo desde 1ra Atención hasta Fallecimiento")
add_subtitle(slide, "De los pacientes que fallecieron, ¿cuánto tiempo pasó desde su primera atención en FISSAL?")
add_footer(slide)

add_kpi_card(slide, "Fallecidos", f"{N_FALL:,} ({PCT_FALL:.1f}%)", 0.5, 1.8, color=BAD)
add_kpi_card(slide, "Sin registro de fallecimiento", f"{N_VIVOS:,} ({100-PCT_FALL:.1f}%)", 3.6, 1.8, color=GOOD)

fallecidos = perfil[perfil["FALLECIDO"]]
if len(fallecidos) > 0:
    sup = fallecidos["SUPERVIVENCIA_DIAS"].dropna()
    sup_meses = sup / 30
    add_kpi_card(slide, "Tiempo mediano hasta muerte", f"{sup_meses.median():.0f} meses", 6.7, 1.8, color=PRIMARY)
    add_kpi_card(slide, "Tiempo medio hasta muerte", f"{sup_meses.mean():.0f} meses", 9.8, 1.8, color=MID_BLUE)

    rangos_sup = [(0, 3), (3, 6), (6, 12), (12, 24), (24, 36), (36, 60), (60, 999)]
    etiquetas_sup = ["0-3 meses", "3-6 meses", "6-12 meses", "1-2 años", "2-3 años", "3-5 años", "5+ años"]
    rows_sup = []
    for (lo, hi), etq in zip(rangos_sup, etiquetas_sup):
        n = ((sup_meses >= lo) & (sup_meses < hi)).sum()
        rows_sup.append([etq, n, f"{n/len(sup)*100:.1f}%"])

    add_text_box(slide, "¿Cuánto tiempo pasó desde su 1ra atención hasta el fallecimiento?", 0.5, 3.6, 8, 0.4, font_size=13, bold=True)
    add_table(slide, ["Tiempo hasta fallecer", "Pacientes", "%"], rows_sup, 0.5, 4.1, [3.0, 1.5, 1.0], font_size=10)

    add_explanation_box(slide, "IMPORTANTE: Esta tabla NO mide 'supervivencia' en el sentido de Kaplan-Meier. Solo describe cuánto tiempo pasó entre la 1ra atención en FISSAL y la muerte, para los pacientes que SÍ fallecieron. No incluye a los pacientes que siguen vivos (censurados). La mediana de Kaplan-Meier usa TODOS los pacientes (vivos y fallecidos).", 0.5, 5.5, 12, 0.9)

# =====================================================================
# SLIDE 5: PERFIL DE TRATAMIENTO
# =====================================================================
slide = add_slide(prs)
add_title(slide, "Tratamientos Recibidos")
add_subtitle(slide, "Indicadores por paciente: ¿qué tipo de atención tuvo al menos una vez?")
add_footer(slide)

tratamientos = [
    ("TUVO_CIRUGIA", "Cirugía", "Proxy: COLECTOM, RESECC, CIRUG, etc. en descripción"),
    ("TUVO_QUIMIOTERAPIA", "Quimioterapia", "Proxy: QUIMIO, OXALIPLAT, FLUOROURAC, etc."),
    ("TUVO_HOSPITALIZACION", "Hospitalización", "Al menos un TIPO_ATENC que contiene HOSP"),
    ("TUVO_EMERGENCIA", "Emergencia", "Al menos un TIPO_ATENC que contiene EMERG"),
    ("TUVO_MEDICAMENTO", "Medicamentos", "Al menos un registro tipo Medicamento"),
    ("TUVO_PROCEDIMIENTO", "Procedimientos", "Al menos un registro tipo Procedimiento"),
]

colors_t = [PRIMARY, MID_BLUE, NEUTRAL, BAD, GOOD, DARK_BLUE]
for i, (col, label, expl) in enumerate(tratamientos):
    n = perfil[col].sum()
    pct = n / N_PAC * 100
    row = i // 3; col_idx = i % 3
    add_kpi_card(slide, label, f"{n:,} ({pct:.1f}%)", 0.5 + col_idx * 4.15, 1.8 + row * 2.1, width=3.8, color=colors_t[i])
    add_text_box(slide, expl, 0.7 + col_idx * 4.15, 3.1 + row * 2.1, 3.5, 0.3, font_size=7, color=GRAY)

add_explanation_box(slide, "Cirugía y Quimioterapia se detectan por palabras clave en la descripción de la prestación (proxies), no por códigos estandarizados de procedimiento. Los porcentajes deben interpretarse como aproximaciones.", 0.5, 6.0, 12, 0.7)

# =====================================================================
# SLIDE 6: INTENSIDAD DE ATENCION
# =====================================================================
slide = add_slide(prs)
add_title(slide, "Intensidad de Atención")
add_subtitle(slide, "¿Cuántas veces se atendió cada paciente y en cuántos establecimientos?")
add_footer(slide)

at = perfil["N_ATENCIONES"]
add_kpi_card(slide, "Atenciones (FUA)\npor paciente — media", f"{at.mean():.1f}", 0.5, 1.8, color=PRIMARY)
add_kpi_card(slide, "Atenciones (FUA)\nmediana", f"{at.median():.0f}", 3.6, 1.8, color=MID_BLUE)
add_kpi_card(slide, "Percentil 25", f"{at.quantile(0.25):.0f}", 6.7, 1.8, color=GOOD)
add_kpi_card(slide, "Percentil 95", f"{at.quantile(0.95):.0f}", 9.8, 1.8, color=BAD)

add_kpi_card(slide, "IPRESS visitados\npor paciente — media", f"{perfil['N_IPRESS'].mean():.1f}", 0.5, 3.4, color=PRIMARY)
add_kpi_card(slide, "Diagnósticos CIE-10\ndistintos — media", f"{perfil['N_DIAGNOSTICOS'].mean():.1f}", 3.6, 3.4, color=MID_BLUE)

tiempo_vivos = perfil.loc[~perfil["FALLECIDO"], "TIEMPO_EN_SISTEMA_DIAS"]
add_kpi_card(slide, "Tiempo en sistema\nde vivos — mediana", f"{tiempo_vivos.median()/30:.0f} meses", 6.7, 3.4, color=NEUTRAL)
add_kpi_card(slide, "Tiempo en sistema\nde vivos — media", f"{tiempo_vivos.mean()/30:.1f} meses", 9.8, 3.4, color=DARK_BLUE)

add_text_box(slide, "¿Qué significa cada métrica?", 0.5, 5.2, 12, 0.3, font_size=13, bold=True)
add_bullet(slide, "FUA = Formato Único de Atención. Cada FUA es una visita o episodio de atención. Un paciente puede tener múltiples prestaciones dentro de una misma FUA.", 0.5, 5.6, 12, 0.4, font_size=10, color=GRAY)
add_bullet(slide, "IPRESS = Instituciones Prestadoras de Servicios de Salud. Cuántas IPRESS distintas visitó el paciente en todo su seguimiento.", 0.5, 6.0, 12, 0.4, font_size=10, color=GRAY)
add_bullet(slide, "Tiempo en sistema = días entre la primera y última atención registrada en FISSAL. Para pacientes vivos, indica cuánto tiempo llevan en seguimiento.", 0.5, 6.4, 12, 0.4, font_size=10, color=GRAY)

# =====================================================================
# SLIDE 7: ENFERMEDAD AVANZADA
# =====================================================================
slide = add_slide(prs)
add_title(slide, "Proxies de Enfermedad Avanzada / Metastásica")
add_subtitle(slide, "FISSAL no registra metástasis (C77/C79). Se usan indicadores indirectos.")
add_footer(slide)

if "ENFERMEDAD_AVANZADA" in perfil_amp.columns:
    avanzada = perfil_amp["ENFERMEDAD_AVANZADA"].value_counts()
    orden = ["PALIATIVO", "MULTISITIO", "UN_SITIO", "SIN_INDICADORES"]
    colors_av = [BAD, NEUTRAL, MID_BLUE, GOOD]
    labels_av = ["Paliativo", "Multisitio (≥2)", "Un sitio", "Sin indicadores"]
    for i, (cat, label) in enumerate(zip(orden, labels_av)):
        n = avanzada.get(cat, 0)
        pct = n / N_PAC * 100
        add_kpi_card(slide, label, f"{n:,} ({pct:.2f}%)", 0.5 + i * 3.15, 1.8, width=2.9, color=colors_av[i])

    sitios = ["SITIO_PERITONEO", "SITIO_PULMON", "SITIO_GANGLIO", "SITIO_MEDULA_OSEA", "SITIO_HIGADO"]
    labels_sitios = ["Peritoneo", "Pulmón", "Ganglio linfático", "Médula ósea", "Hígado"]
    rows_sitios = []
    for col, lbl in zip(sitios, labels_sitios):
        n = perfil_amp[col].sum() if col in perfil_amp.columns else 0
        rows_sitios.append([lbl, n])

    add_text_box(slide, "Procedimientos en órganos a distancia (solo procedimientos, excluye insumos/medicamentos)", 0.5, 3.6, 8, 0.4, font_size=13, bold=True)
    add_table(slide, ["Órgano / Sitio", "Pacientes"], rows_sitios, 0.5, 4.1, [3.0, 1.5], font_size=10)

    add_explanation_box(slide, "Paliativo = visita domiciliaria de cuidados paliativos (26 pacientes, solo 2016-2018).  Multisitio = procedimientos en ≥2 órganos a distancia.  Un Sitio = procedimiento en 1 órgano.  Estos son los únicos proxies disponibles en FISSAL para inferir enfermedad avanzada.", 0.5, 5.3, 12, 0.8)

    add_text_box(slide, "Perfil de cada grupo: Paliativo: 26 pac, 26.9% fall, 42.3% cirugía, 42.3% quimio  |  Un Sitio: 126 pac, 16.7% fall, 79.4% cirugía, 81% hospitalización  |  Sin indicadores: 15,544 pac, 15.2% fall", 0.5, 6.3, 12, 0.5, font_size=9, color=GRAY)

# =====================================================================
# SLIDE 8: COMO INTERPRETAR KAPLAN-MEIER
# =====================================================================
slide = add_slide(prs)
add_title(slide, "¿Cómo interpretar las curvas de Kaplan-Meier?")
add_subtitle(slide, "Guía rápida para leer los gráficos de supervivencia")
add_footer(slide)

explicaciones = [
    ("Eje X: tiempo", "Meses transcurridos desde la primera atención en FISSAL."),
    ("Eje Y: probabilidad", "Probabilidad de seguir vivo. Empieza en 1.0 (100%) y va bajando con cada fallecimiento."),
    ("Escalones hacia abajo", "Cada vez que la curva baja, uno o más pacientes fallecieron en ese momento."),
    ("Marcas de censura (+)", "Pacientes que siguen vivos al final del estudio o que se perdieron. No bajan la curva."),
    ("Mediana de supervivencia", "Tiempo donde la curva cruza el 0.5 (50%). Si nunca cruza, la mayoría sigue viva al corte."),
    ("Bandas sombreadas", "Intervalo de confianza del 95%. Si dos bandas no se solapan, la diferencia es significativa."),
    ("Log-rank test (p-value)", "Prueba estadística: si p < 0.05, hay diferencia significativa entre las curvas."),
]

for i, (tit, desc) in enumerate(explicaciones):
    y = 1.7 + i * 0.75
    add_bullet(slide, f"▪ {tit}", 0.7, y, 5, 0.35, font_size=12, bold=True, color=PRIMARY)
    add_text_box(slide, desc, 1.0, y + 0.3, 11, 0.35, font_size=10, color=FOREGROUND)

add_explanation_box(slide, "IMPORTANTE: En este estudio, la mediana de supervivencia NO se alcanza para la mayoría de grupos porque solo el 15.2% de pacientes falleció durante el seguimiento. Las curvas son válidas para comparar grupos entre sí, pero la mediana 'N/A' no significa que no haya muertes.", 0.5, 6.2, 12, 0.7)

# =====================================================================
# SLIDES 9-14: KAPLAN-MEIER individuales + SLIDE 15 RESUMEN
# =====================================================================
km_data = [
    ("km_global.png", "Supervivencia Global", "Todo paciente con cáncer colorrectal (C18/C19/C20). Línea base de referencia."),
    ("km_sexo.png", "Supervivencia por Sexo", "Femenino (F) vs Masculino (M). ¿Hay diferencia de supervivencia entre sexos?"),
    ("km_edad.png", "Supervivencia por Rango de Edad", "¿Los pacientes más jóvenes sobreviven más que los mayores? Estratificación por grupo etario."),
    ("km_cirugia.png", "Supervivencia según Cirugía", "Pacientes que recibieron cirugía vs los que no. La cirugía es el tratamiento curativo principal en CCR."),
    ("km_quimio.png", "Supervivencia según Quimioterapia", "Pacientes con vs sin quimioterapia. Nota: posible sesgo de selección (reciben quimio los que están en condiciones de recibirla)."),
    ("km_enfermedad_avanzada.png", "Supervivencia según Enfermedad Avanzada", "Proxy compuesto: paliativo, multisitio (≥2 órganos), un sitio, sin indicadores. Grupos pequeños = bandas de confianza anchas."),
    ("km_region.png", "Supervivencia: Lima vs Regiones", "¿Hay diferencias en supervivencia entre Lima Metropolitana y el resto del país?"),
]

for filename, titulo, descripcion in km_data:
    slide = add_slide(prs)
    add_title(slide, titulo)
    add_subtitle(slide, descripcion)
    add_footer(slide)

    img_path = OUTPUT_DIR / filename
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.0))
    else:
        add_text_box(slide, f"[Gráfico no encontrado: {filename}]", 1, 2.5, 10, 1, font_size=16, color=BAD)

# SLIDE 15: resumen KM
slide = add_slide(prs)
add_title(slide, "Resumen Comparativo — 6 Dimensiones")
add_subtitle(slide, "Visión conjunta de todas las curvas Kaplan-Meier")
add_footer(slide)
img_path = OUTPUT_DIR / "km_resumen_comparativo.png"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.5))

# =====================================================================
# SLIDE 16: TIEMPOS ENTRE HITOS CLINICOS
# =====================================================================
print("trayectoria...", end=" ")
slide = add_slide(prs)
add_title(slide, "Tiempos entre Hitos Clínicos")
add_subtitle(slide, "Días desde la primera atención hasta cada evento clave del tratamiento")
add_footer(slide)

# Calcular hitos a partir de df_ccr
print("calculando hitos...", end=" ")
PAT_CAT = {
    "BIOPSIA": r"(?i)BIOPSI|ANATOMOPATOLOG",
    "ENDOSCOPIA": r"(?i)COLONOSCOP|ENDOSCOP|RECTOSCOP|SIGMOIDOSCOP",
    "CIRUGIA": r"(?i)COLECTOM|RESECC|HEMICOLEC|COLOSTOM|LAPAROTOM|ANASTOMOS|LAPAROSCOP",
    "IMAGEN": r"(?i)TOMOGRAF|RESONAN|ECOGRAF|GAMMAGRAF|PET.?SCAN",
    "QUIMIOTERAPIA": r"(?i)QUIMIOTER|OXALIPLAT|FLUOROURAC|CAPECITAB|IRINOTECAN|BEVACIZUM|CETUXIMAB|FOLFOX|FOLFIRI|LEUCOVORIN",
    "RADIOTERAPIA": r"(?i)RADIOTER|COBALTOT|BRAQUIT|ACELERADOR",
}

primera_at = df_ccr.groupby("NUM_DOC_CRYPTO")["ATE_FECATENCION"].min()

hitos_rows = []
for nombre, patron in PAT_CAT.items():
    mask = df_ccr["ATE_DESCCONSUMO"].str.contains(patron, na=False)
    df_cat = df_ccr[mask]
    if len(df_cat) == 0:
        continue
    primera_cat = df_cat.groupby("NUM_DOC_CRYPTO")["ATE_FECATENCION"].min()
    merged = pd.merge(primera_cat.rename("FECHA"), primera_at.rename("FECHA_1RA"), left_index=True, right_index=True)
    dias = (merged["FECHA"] - merged["FECHA_1RA"]).dt.days
    dias_pos = dias[dias >= 0]
    n_pac = len(dias_pos)
    if n_pac > 0:
        hitos_rows.append([
            nombre, n_pac,
            f"{dias_pos.median():.0f}",
            f"{dias_pos.mean():.0f}",
            f"{dias_pos.quantile(0.25):.0f}",
            f"{dias_pos.quantile(0.75):.0f}",
        ])

add_text_box(slide, "Días desde la 1ra atención en FISSAL hasta el 1er procedimiento de cada tipo", 0.5, 1.5, 8, 0.4, font_size=13, bold=True)
add_table(slide, ["Hito clínico", "N pacientes", "Mediana (d)", "Media (d)", "P25 (d)", "P75 (d)"],
          hitos_rows, 0.5, 2.0, [2.5, 1.5, 1.3, 1.3, 1.2, 1.2], font_size=10)

# Tiempo cirugia → quimio
mask_cir = df_ccr["ATE_DESCCONSUMO"].str.contains(PAT_CAT["CIRUGIA"], na=False)
mask_qt = df_ccr["ATE_DESCCONSUMO"].str.contains(PAT_CAT["QUIMIOTERAPIA"], na=False)
primera_cir = df_ccr[mask_cir].groupby("NUM_DOC_CRYPTO")["ATE_FECATENCION"].min()
primera_qt = df_ccr[mask_qt].groupby("NUM_DOC_CRYPTO")["ATE_FECATENCION"].min()
ambos = pd.merge(primera_cir.rename("CIR"), primera_qt.rename("QT"), left_index=True, right_index=True)
dias_cir_qt = (ambos["QT"] - ambos["CIR"]).dt.days
qt_post = dias_cir_qt[dias_cir_qt > 0]
qt_pre = dias_cir_qt[dias_cir_qt < 0]

add_text_box(slide, "Tiempo entre Cirugía y Quimioterapia", 0.5, 4.5, 8, 0.4, font_size=13, bold=True)
add_text_box(slide, f"Pacientes con ambos: {len(ambos):,}  |  Quimio DESPUÉS de cirugía (adyuvante): {len(qt_post):,}  |  Quimio ANTES de cirugía (neoadyuvante): {len(qt_pre):,}", 0.5, 5.0, 12, 0.3, font_size=11)
if len(qt_post) > 0:
    add_text_box(slide, f"Adyuvante (cirugía → quimio): mediana = {qt_post.median():.0f} días ({qt_post.median()/30:.1f} meses), media = {qt_post.mean():.0f} días ({qt_post.mean()/30:.1f} meses)", 0.5, 5.4, 12, 0.3, font_size=11)
if len(qt_pre) > 0:
    add_text_box(slide, f"Neoadyuvante (quimio → cirugía): mediana = {abs(qt_pre).median():.0f} días ({abs(qt_pre).median()/30:.1f} meses), media = {abs(qt_pre).mean():.0f} días ({abs(qt_pre).mean()/30:.1f} meses)", 0.5, 5.8, 12, 0.3, font_size=11)

add_explanation_box(slide, "Endoscopia y Biopsia deberían ser los primeros hitos (diagnóstico). Cirugía y Quimioterapia vienen después (tratamiento). La mediana muestra el tiempo 'típico' hasta cada hito.", 0.5, 6.3, 12, 0.5)

# =====================================================================
# SLIDE 17: ESQUEMAS DE QUIMIOTERAPIA
# =====================================================================
print("quimio...", end=" ")
slide = add_slide(prs)
add_title(slide, "Esquemas de Quimioterapia")
add_subtitle(slide, "Regímenes identificados por combinación de drogas en descripciones de medicamentos")
add_footer(slide)

med = df_ccr[df_ccr["ATE_TIPOCONSUMO"] == "Medicamento"].copy()
PAT_QUIMIO_DRUGS = {
    "OXALIPLATINO": r"(?i)OXALIPLAT",
    "FLUOROURACILO": r"(?i)FLUOROURAC",
    "CAPECITABINA": r"(?i)CAPECITAB",
    "IRINOTECAN": r"(?i)IRINOTECAN",
    "LEUCOVORIN": r"(?i)FOLINATO|LEUCOVORIN",
    "BEVACIZUMAB": r"(?i)BEVACIZUM",
    "CETUXIMAB": r"(?i)CETUXIMAB",
}

qt_drogas = {}
for drug, pat in PAT_QUIMIO_DRUGS.items():
    mask = med["ATE_DESCCONSUMO"].str.contains(pat, na=False)
    pacs = med.loc[mask, "NUM_DOC_CRYPTO"].unique()
    for p in pacs:
        qt_drogas.setdefault(p, set()).add(drug)

qt_drogas_set = {p: frozenset(d) for p, d in qt_drogas.items()}

PAT_ESQUEMA = {
    "FOLFOX": lambda d: {"FLUOROURACILO", "LEUCOVORIN", "OXALIPLATINO"}.issubset(d),
    "CAPOX / XELOX": lambda d: {"CAPECITABINA", "OXALIPLATINO"}.issubset(d) and "FLUOROURACILO" not in d,
    "FOLFIRI": lambda d: {"FLUOROURACILO", "LEUCOVORIN", "IRINOTECAN"}.issubset(d),
    "CAPECITABINA mono": lambda d: d == {"CAPECITABINA"},
    "5FU / LV": lambda d: {"FLUOROURACILO", "LEUCOVORIN"}.issubset(d) and len(d) == 2,
    "5FU solo": lambda d: d == {"FLUOROURACILO"},
    "OXALIPLATINO solo": lambda d: d == {"OXALIPLATINO"},
}

regimenes = {p: "Otro / combinación" for p in qt_drogas_set}
for nombre, fn in PAT_ESQUEMA.items():
    for p, d in qt_drogas_set.items():
        if fn(d):
            regimenes[p] = nombre

reg_counts = pd.Series(regimenes).value_counts()
n_beva = sum(1 for d in qt_drogas_set.values() if "BEVACIZUMAB" in d)
n_cetux = sum(1 for d in qt_drogas_set.values() if "CETUXIMAB" in d)
n_total_qt = len(qt_drogas_set)

rows_reg = [[reg, n, f"{n/n_total_qt*100:.1f}%"] for reg, n in reg_counts.items()]
add_text_box(slide, f"Pacientes con al menos una droga de quimioterapia identificada: {n_total_qt:,}", 0.5, 1.5, 8, 0.4, font_size=13, bold=True)
add_table(slide, ["Régimen", "Pacientes", "%"], rows_reg, 0.5, 2.0, [3.0, 1.5, 1.0], font_size=10)

add_text_box(slide, f"Pacientes con Bevacizumab añadido: {n_beva:,}  |  Pacientes con Cetuximab añadido: {n_cetux:,}", 0.5, 5.2, 12, 0.4, font_size=12, bold=True)

add_explanation_box(slide, "Los regímenes se infieren por las drogas que aparecen en los registros de medicamentos del paciente. FOLFOX y CAPOX son los estándar para cáncer de colon. Bevacizumab y Cetuximab son terapias dirigidas (anticuerpos monoclonales) que se añaden a la quimioterapia en enfermedad metastásica.", 0.5, 5.8, 12, 0.8)

# =====================================================================
# SLIDE 18: COSTOS
# =====================================================================
print("costos...", end=" ")
slide = add_slide(prs)
add_title(slide, "Costos de la Atención")
add_subtitle(slide, "Distribución del monto neto total por paciente y por tipo de consumo")
add_footer(slide)

costo_total = perfil["MONTO_NETO_TOTAL"]
add_kpi_card(slide, "Costo neto total\n(toda la cohorte)", f"S/ {costo_total.sum():,.0f}", 0.5, 1.8, color=PRIMARY)
add_kpi_card(slide, "Costo medio\npor paciente", f"S/ {costo_total.mean():,.0f}", 3.6, 1.8, color=MID_BLUE)
add_kpi_card(slide, "Costo mediano\npor paciente", f"S/ {costo_total.median():,.0f}", 6.7, 1.8, color=GOOD)
add_kpi_card(slide, "Percentil 95", f"S/ {costo_total.quantile(0.95):,.0f}", 9.8, 1.8, color=BAD)

# Costo por tipo de consumo
costos_tipo = df_ccr.groupby("ATE_TIPOCONSUMO")["ATE_MONTONETO"].sum().sort_values(ascending=False)
rows_costo = []
for tipo, monto in costos_tipo.items():
    pct = monto / costo_total.sum() * 100
    rows_costo.append([tipo, f"S/ {monto:,.0f}", f"{pct:.1f}%"])

add_text_box(slide, "Costo neto total por tipo de consumo", 0.5, 3.5, 6, 0.4, font_size=13, bold=True)
add_table(slide, ["Tipo de consumo", "Monto neto", "%"], rows_costo, 0.5, 4.0, [2.5, 2.0, 1.0], font_size=10)

# Costo por año
costos_anio = df_ccr.groupby("ANIO_ATENCION")["ATE_MONTONETO"].sum()
rows_ca = [[int(yr), f"S/ {monto:,.0f}"] for yr, monto in costos_anio.items()]
add_text_box(slide, "Costo neto por año de atención", 7.5, 3.5, 5, 0.4, font_size=13, bold=True)
add_table(slide, ["Año", "Monto neto"], rows_ca, 7.5, 4.0, [1.5, 2.0], font_size=10)

add_explanation_box(slide, "El costo está concentrado en medicamentos (quimioterapia, terapia dirigida) y procedimientos (cirugías). La distribución es asimétrica: unos pocos pacientes concentran gran parte del gasto (ver P95 vs mediana).", 0.5, 6.2, 12, 0.6)

# =====================================================================
# SLIDE 19: HOSPITALIZACION
# =====================================================================
print("hosp...", end=" ")
slide = add_slide(prs)
add_title(slide, "Hospitalizaciones")
add_subtitle(slide, "Episodios de hospitalización, días de estancia y distribución")
add_footer(slide)

hosp = df_ccr[df_ccr["TIPO_ATENC"].str.contains("HOSP", case=False, na=False)].copy()
n_episodios = len(hosp)
n_pac_hosp = hosp["NUM_DOC_CRYPTO"].nunique()

if "ATE_FECINGHOSP" in hosp.columns and "ATE_FECALTHOSP" in hosp.columns:
    hosp["DIAS_ESTANCIA"] = (pd.to_datetime(hosp["ATE_FECALTHOSP"]) - pd.to_datetime(hosp["ATE_FECINGHOSP"])).dt.days
    estancia = hosp["DIAS_ESTANCIA"].dropna()
    estancia = estancia[(estancia >= 0) & (estancia <= 365)]
else:
    estancia = pd.Series(dtype=float)

add_kpi_card(slide, "Episodios de\nhospitalización", f"{n_episodios:,}", 0.5, 1.8, color=PRIMARY)
add_kpi_card(slide, "Pacientes\nhospitalizados", f"{n_pac_hosp:,} ({n_pac_hosp/N_PAC*100:.1f}%)", 3.6, 1.8, color=MID_BLUE)
if len(estancia) > 0:
    add_kpi_card(slide, "Estancia mediana", f"{estancia.median():.0f} días", 6.7, 1.8, color=GOOD)
    add_kpi_card(slide, "Estancia media", f"{estancia.mean():.1f} días", 9.8, 1.8, color=BAD)
else:
    add_kpi_card(slide, "Estancia mediana", "N/D", 6.7, 1.8, color=GRAY)
    add_kpi_card(slide, "Estancia media", "N/D", 9.8, 1.8, color=GRAY)

# Hospitalizaciones por año
hosp_anio = hosp.groupby("ANIO_ATENCION")["NUM_DOC_CRYPTO"].nunique()
rows_ha = [[int(yr), n] for yr, n in hosp_anio.items()]
add_text_box(slide, "Pacientes hospitalizados por año", 0.5, 3.5, 5, 0.4, font_size=13, bold=True)
add_table(slide, ["Año", "Pacientes"], rows_ha, 0.5, 4.0, [1.5, 1.5], font_size=10)

if len(estancia) > 0:
    bins_est = [0, 1, 3, 7, 15, 30, 365]
    labels_est = ["1 día", "2-3 días", "4-7 días", "8-15 días", "16-30 días", "30+ días"]
    est_rango = pd.cut(estancia, bins=bins_est, labels=labels_est, right=True)
    rows_est = []
    for lbl in labels_est:
        n = (est_rango == lbl).sum()
        rows_est.append([lbl, n, f"{n/len(estancia)*100:.1f}%"])
    add_text_box(slide, "Distribución de días de estancia", 6.8, 3.5, 5, 0.4, font_size=13, bold=True)
    add_table(slide, ["Días de estancia", "Episodios", "%"], rows_est, 6.8, 4.0, [2.0, 1.5, 1.0], font_size=10)

# =====================================================================
# SLIDE 20: RECURRENCIAS (GAPS)
# =====================================================================
print("recurrencias...", end=" ")
slide = add_slide(prs)
add_title(slide, "Brechas en la Atención (Proxy de Recurrencia)")
add_subtitle(slide, "Pacientes con gaps > 180 días sin atención, posible indicador de recurrencia o abandono")
add_footer(slide)

df_ord = df_ccr.sort_values(["NUM_DOC_CRYPTO", "ATE_FECATENCION"])
df_ord["FECHA_ANTERIOR"] = df_ord.groupby("NUM_DOC_CRYPTO")["ATE_FECATENCION"].shift(1)
df_ord["DIAS_GAP"] = (df_ord["ATE_FECATENCION"] - df_ord["FECHA_ANTERIOR"]).dt.days
gaps = df_ord[df_ord["DIAS_GAP"] > 180]["DIAS_GAP"]

n_pac_gap = df_ord.loc[df_ord["DIAS_GAP"] > 180, "NUM_DOC_CRYPTO"].nunique()
n_ep_gap = len(gaps)

add_kpi_card(slide, "Episodios de\ngap > 180 días", f"{n_ep_gap:,}", 0.5, 1.8, color=NEUTRAL)
add_kpi_card(slide, "Pacientes con\nal menos 1 gap", f"{n_pac_gap:,} ({n_pac_gap/N_PAC*100:.1f}%)", 3.6, 1.8, color=MID_BLUE)
add_kpi_card(slide, "Gap mediano", f"{gaps.median():.0f} días" if len(gaps) > 0 else "N/D", 6.7, 1.8, color=GOOD)
add_kpi_card(slide, "Gap máximo", f"{gaps.max():.0f} días" if len(gaps) > 0 else "N/D", 9.8, 1.8, color=BAD)

if len(gaps) > 0:
    rangos_gap = [(180, 365), (365, 730), (730, 1095), (1095, 9999)]
    etq_gap = ["6-12 meses", "1-2 años", "2-3 años", "3+ años"]
    rows_gap = []
    for (lo, hi), etq in zip(rangos_gap, etq_gap):
        n = ((gaps >= lo) & (gaps < hi)).sum()
        rows_gap.append([etq, n, f"{n/len(gaps)*100:.1f}%"])
    add_text_box(slide, "Distribución de la duración de los gaps", 0.5, 3.5, 6, 0.4, font_size=13, bold=True)
    add_table(slide, ["Duración del gap", "Episodios", "%"], rows_gap, 0.5, 4.0, [2.5, 1.5, 1.0], font_size=10)

add_explanation_box(slide, "Un gap > 180 días sin atenciones puede indicar: (a) recurrencia tratada en otro establecimiento fuera de FISSAL, (b) abandono de tratamiento, (c) el paciente fue dado de alta y luego recayó, o (d) periodo sin necesidad de atención. No se puede distinguir entre estas causas con los datos disponibles.", 0.5, 5.0, 12, 1.0)

del df_ord, df_ccr; gc.collect()

# =====================================================================
# SLIDE 21: LIMITACIONES
# =====================================================================
slide = add_slide(prs)
add_title(slide, "Limitaciones y Notas Metodológicas")
add_subtitle(slide, "Consideraciones importantes para interpretar los resultados")
add_footer(slide)

limitaciones = [
    ("Sin códigos de metástasis (C77/C79)", "FISSAL no registra códigos de neoplasia secundaria. La detección de enfermedad avanzada usa proxies indirectos (paliativos, procedimientos en otros órganos)."),
    (f"Mediana de supervivencia no alcanzada", f"Solo {PCT_FALL:.1f}% falleció en el periodo de estudio. La mayoría sigue viva al corte (31-dic-2022). Las curvas KM nunca cruzan el 50%."),
    ("Sin estadiaje TNM", "No hay datos de estadio clínico en FISSAL. No se puede diferenciar entre estadios tempranos y avanzados."),
    ("Paliativos solo 2016-2018", "El registro 'Visita domiciliaria de seguimiento - Cuidados Paliativos' desaparece después de 2018. Solo 26 pacientes."),
    ("Cirugía y quimioterapia: aproximaciones", "Se detectan por palabras clave en descripciones de consumo, no por códigos estandarizados de procedimiento (CPT)."),
    ("Periodo de seguimiento variable", "Aunque hay atenciones desde 2012, la mayoría de datos se concentra en 2016-2022. El seguimiento máximo real es de ~10 años pero para pocos pacientes."),
    ("Sesgo de selección en quimioterapia", "Los pacientes que reciben quimio pueden tener mejor estado general, lo que sesga la comparación de supervivencia."),
    ("Sin causa de muerte", "FISSAL registra fecha de fallecimiento pero no la causa. No se sabe si la muerte fue por cáncer u otra razón."),
]

y = 1.6
for tit, desc in limitaciones:
    add_bullet(slide, f"▪ {tit}", 0.7, y, 11.5, 0.35, font_size=12, bold=True, color=BAD)
    add_text_box(slide, desc, 1.1, y + 0.3, 11.3, 0.4, font_size=9, color=GRAY)
    y += 0.65

# =====================================================================
# GUARDAR
# =====================================================================
PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(PPTX_OUT))
print(f"\nPresentacion guardada: {PPTX_OUT}")
print(f"Total slides: {len(prs.slides)}")
